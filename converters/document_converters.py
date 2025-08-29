import os
import uuid
import shutil
import tempfile
import subprocess
from pathlib import Path
from typing import Optional

# Lazy imports inside functions to keep cross-platform compatibility

def _html_to_pdf_text_fallback(html_path: str, out_path: str):
    """Fallback: render HTML as plain text in PDF (no formatting, just text)."""
    from reportlab.pdfgen import canvas
    from reportlab.lib.pagesizes import letter
    _ensure_dir(Path(out_path).parent.as_posix())
    with open(html_path, 'r', encoding='utf-8') as f:
        lines = f.readlines()
    c = canvas.Canvas(out_path, pagesize=letter)
    y = 750
    for line in lines:
        if y < 50:
            c.showPage()
            y = 750
        c.setFont("Helvetica", 12)
        c.drawString(50, y, line.rstrip()[:100])
        y -= 14
    c.save()
    return out_path

def _remove_excel_metadata(doc_path: str, out_path: str):
    from openpyxl import load_workbook
    wb = load_workbook(doc_path)
    props = wb.properties
    props.creator = ""
    props.lastModifiedBy = ""
    props.created = None
    props.modified = None
    props.title = ""
    props.subject = ""
    props.description = ""
    props.keywords = ""
    props.category = ""
    wb.save(out_path)
    return out_path
# converters/document_converters.py


# ============== Utils ==============

def _ensure_dir(path: str):
    os.makedirs(path, exist_ok=True)

def _unique_name(prefix: str, ext: str) -> str:
    return f"{prefix}_{uuid.uuid4().hex}{ext}"

def _which(candidates):
    from shutil import which
    for c in candidates:
        p = which(c)
        if p:
            return p
    return None


# ============== PDF <-> Word ==============

def pdf_to_docx(pdf_path: str, out_path: str):
    """
    Convert PDF -> DOCX using pdf2docx.
    Note: Complex PDFs may not convert perfectly.
    """
    _ensure_dir(Path(out_path).parent.as_posix())
    from pdf2docx import Converter
    cv = Converter(pdf_path)
    try:
        cv.convert(out_path, start=0, end=None)
    finally:
        cv.close()
    return out_path

def html_to_pdf(html_path: str, out_path: str):
    """Convert HTML to PDF using multiple engines"""
    _ensure_dir(Path(out_path).parent.as_posix())
    # Method 1: WeasyPrint (best for CSS)
    try:
        import weasyprint
        weasyprint.HTML(filename=html_path).write_pdf(out_path)
        return out_path
    except ImportError:
        pass
    # Method 2: pdfkit + wkhtmltopdf
    try:
        import pdfkit
        pdfkit.from_file(html_path, out_path)
        return out_path
    except Exception:
        pass
    # Method 3: ReportLab text fallback
    return _html_to_pdf_text_fallback(html_path, out_path)


def _docx_to_pdf_windows_com(docx_path: str, out_path: str):
    """
    Windows: Use Word COM automation for best fidelity.
    Requires Microsoft Word installed.
    """
    import win32com.client  # pip install pywin32
    wdFormatPDF = 17
    word = win32com.client.DispatchEx("Word.Application")
    word.Visible = False
    doc = word.Documents.Open(os.path.abspath(docx_path))
    try:
        doc.ExportAsFixedFormat(os.path.abspath(out_path), wdFormatPDF)
    finally:
        doc.Close(False)
        word.Quit()
    return out_path


def _docx_to_pdf_libreoffice(docx_path: str, out_path: str):
    """
    Cross-platform fallback using LibreOffice (soffice) headless.
    """
    soffice = _which(["soffice", "libreoffice"])
    if not soffice:
        raise RuntimeError("LibreOffice (soffice) not found in PATH.")
    out_dir = str(Path(out_path).parent)
    # Convert; LibreOffice writes file as <name>.pdf in out_dir
    subprocess.run([soffice, "--headless", "--convert-to", "pdf", "--outdir", out_dir, docx_path], check=True)
    generated = os.path.join(out_dir, Path(docx_path).with_suffix(".pdf").name)
    # If name mismatched (rare), find a PDF and move
    if not os.path.exists(generated):
        # fallback: pick first pdf created
        candidates = [f for f in os.listdir(out_dir) if f.lower().endswith(".pdf")]
        if not candidates:
            raise RuntimeError("LibreOffice did not produce a PDF.")
        generated = os.path.join(out_dir, candidates[0])
    if os.path.abspath(generated) != os.path.abspath(out_path):
        shutil.move(generated, out_path)
    return out_path

def txt_to_pdf(txt_path: str, out_path: str, font_size: int = 12):
    """Convert text file to PDF with formatting"""
    from reportlab.pdfgen import canvas
    from reportlab.lib.pagesizes import letter
    _ensure_dir(Path(out_path).parent.as_posix())
    with open(txt_path, 'r', encoding='utf-8') as f:
        lines = f.readlines()
    c = canvas.Canvas(out_path, pagesize=letter)
    y = 750
    for line in lines:
        if y < 50:
            c.showPage()
            y = 750
        c.setFont("Helvetica", font_size)
        c.drawString(50, y, line.rstrip()[:100])  # Wrap long lines
        y -= font_size + 2
    c.save()
    return out_path


def _docx_to_pdf_text_fallback(docx_path: str, out_path: str):
    """
    Last-resort: extract text via python-docx and write using ReportLab.
    Loses formatting; only plain text preserved.
    """
    from docx import Document
    from reportlab.pdfgen import canvas
    from reportlab.lib.pagesizes import letter

    doc = Document(docx_path)
    _ensure_dir(Path(out_path).parent.as_posix())
    c = canvas.Canvas(out_path, pagesize=letter)
    y = 750
    for p in doc.paragraphs:
        if y < 50:
            c.showPage()
            y = 750
        txt = p.text.replace("\t", "    ")
        c.drawString(50, y, txt[:2000])
        y -= 18
    c.save()
    return out_path

def merge_word_documents(doc_paths: list, out_path: str):
    """Merge multiple Word documents"""
    from docx import Document
    from docx.oxml.table import CT_Tbl
    from docx.oxml.text.paragraph import CT_P
    merged_doc = Document()
    for doc_path in doc_paths:
        sub_doc = Document(doc_path)
        for element in sub_doc.element.body:
            if isinstance(element, CT_P):
                merged_doc.element.body.append(element)
            elif isinstance(element, CT_Tbl):
                merged_doc.element.body.append(element)
        merged_doc.add_page_break()
    merged_doc.save(out_path)
    return out_path

    """
    Last-resort: extract text via python-docx and write using ReportLab.
    Loses formatting; only plain text preserved.
    """
    from docx import Document
    from reportlab.pdfgen import canvas
    from reportlab.lib.pagesizes import letter

    doc = Document(docx_path)
    _ensure_dir(Path(out_path).parent.as_posix())
    c = canvas.Canvas(out_path, pagesize=letter)
    y = 750
    for p in doc.paragraphs:
        if y < 50:
            c.showPage()
            y = 750
        txt = p.text.replace("\t", "    ")
        c.drawString(50, y, txt[:2000])
        y -= 18
    c.save()
    return out_path


# In converters/document_converters.py

def docx_to_pdf(docx_path: str, out_path: str):
    """
    DOC/DOCX -> PDF with multi-tier fallback:
    1) Windows COM (Word) if available
    2) LibreOffice headless
    3) Text-only ReportLab fallback (DOCX only; DOC not supported here)
    """
    _ensure_dir(Path(out_path).parent.as_posix())

    ext = Path(docx_path).suffix.lower()
    is_doc = (ext == ".doc")
    is_docx = (ext == ".docx")

    # Try Windows COM first
    if os.name == "nt":
        try:
            return _docx_to_pdf_windows_com(docx_path, out_path)
        except Exception:
            # Fall back to LibreOffice on Windows if available
            try:
                return _docx_to_pdf_libreoffice(docx_path, out_path)
            except Exception:
                if is_docx:
                    # Only for DOCX we can text-fallback
                    return _docx_to_pdf_text_fallback(docx_path, out_path)
                # For DOC, clearly tell user what to install
                raise RuntimeError("Cannot convert .doc without Microsoft Word or LibreOffice (soffice) installed.")

    # Non-Windows: try LibreOffice
    try:
        return _docx_to_pdf_libreoffice(docx_path, out_path)
    except Exception:
        if is_docx:
            return _docx_to_pdf_text_fallback(docx_path, out_path)
        raise RuntimeError("Cannot convert .doc on this system without LibreOffice (soffice) available in PATH.")

def protect_word_document(doc_path: str, out_path: str, password: str):
    """Password protect Word document"""
    if os.name == "nt":
        try:
            import win32com.client
            word = win32com.client.DispatchEx("Word.Application")
            word.Visible = False
            doc = word.Documents.Open(os.path.abspath(doc_path))
            doc.Password = password
            doc.SaveAs2(os.path.abspath(out_path))
            doc.Close()
            word.Quit()
            return out_path
        except Exception:
            pass
    shutil.copy2(doc_path, out_path)
    raise RuntimeWarning("Password protection requires Microsoft Word")


# ============== PDF -> Excel ==============

def _pdf_tables_to_excel_camelot(pdf_path: str, out_xlsx_path: str):
    """
    Use camelot-py to extract tables into an Excel workbook.
    Requires: camelot-py[cv] + Ghostscript
    """
    import pandas as pd
    _ensure_dir(Path(out_xlsx_path).parent.as_posix())

    try:
        import camelot
    except ImportError:
        raise RuntimeError("camelot-py not installed. pip install camelot-py[cv]")

    # Try lattice first, then stream
    tables = []
    for flavor in ("lattice", "stream"):
        try:
            t = camelot.read_pdf(pdf_path, pages="all", flavor=flavor)
            if t and len(t) > 0:
                tables = t
                break
        except Exception:
            continue

    if not tables:
        raise RuntimeError("No tables detected by Camelot.")

    with pd.ExcelWriter(out_xlsx_path, engine="openpyxl") as writer:
        for idx, table in enumerate(tables, start=1):
            df = table.df
            df.to_excel(writer, index=False, sheet_name=f"Table_{idx}")
    return out_xlsx_path


def _pdf_text_to_excel_basic(pdf_path: str, out_xlsx_path: str):
    """
    Basic fallback: extract text lines into first column (no table structure).
    """
    from openpyxl import Workbook
    try:
        from PyPDF2 import PdfReader
    except Exception:
        raise RuntimeError("PyPDF2 not available for text extraction.")

    reader = PdfReader(pdf_path)
    wb = Workbook()
    ws = wb.active
    ws.title = "PDF Content"

    for page in reader.pages:
        text = page.extract_text() or ""
        for line in text.splitlines():
            ws.append([line])
        ws.append([""])  # blank separator

    _ensure_dir(Path(out_xlsx_path).parent.as_posix())
    wb.save(out_xlsx_path)
    return out_xlsx_path


def pdf_to_excel(pdf_path: str, out_xlsx_path: str):
    """
    PDF -> Excel:
    1) Try tables via Camelot
    2) Fallback to basic text dump
    """
    try:
        return _pdf_tables_to_excel_camelot(pdf_path, out_xlsx_path)
    except Exception:
        return _pdf_text_to_excel_basic(pdf_path, out_xlsx_path)

def redact_text_from_word(doc_path: str, out_path: str, patterns: list):
    """Redact sensitive text from Word document"""
    from docx import Document
    import re
    doc = Document(doc_path)
    for paragraph in doc.paragraphs:
        for pattern in patterns:
            paragraph.text = re.sub(pattern, "█████", paragraph.text, flags=re.IGNORECASE)
    for table in doc.tables:
        for row in table.rows:
            for cell in row.cells:
                for pattern in patterns:
                    cell.text = re.sub(pattern, "█████", cell.text, flags=re.IGNORECASE)
    doc.save(out_path)
    return out_path


# ============== Excel -> PDF ==============

def _excel_to_pdf_windows_com(xlsx_path: str, out_path: str):
    """
    Windows: Use Excel COM ExportAsFixedFormat
    """
    import win32com.client  # pip install pywin32
    xlTypePDF = 0
    excel = win32com.client.DispatchEx("Excel.Application")
    excel.Visible = False
    wb = excel.Workbooks.Open(os.path.abspath(xlsx_path))
    try:
        wb.ExportAsFixedFormat(xlTypePDF, os.path.abspath(out_path))
    finally:
        wb.Close(False)
        excel.Quit()
    return out_path


def _office_to_pdf_libreoffice_any(in_path: str, out_path: str):
    """
    LibreOffice headless converter for Office docs (Excel, Word, PowerPoint)
    """
    soffice = _which(["soffice", "libreoffice"])
    if not soffice:
        raise RuntimeError("LibreOffice (soffice) not found in PATH.")
    out_dir = str(Path(out_path).parent)
    subprocess.run([soffice, "--headless", "--convert-to", "pdf", "--outdir", out_dir, in_path], check=True)
    generated = os.path.join(out_dir, Path(in_path).with_suffix(".pdf").name)
    if not os.path.exists(generated):
        # Fallback: pick first pdf created
        candidates = [f for f in os.listdir(out_dir) if f.lower().endswith(".pdf")]
        if not candidates:
            raise RuntimeError("LibreOffice did not produce a PDF.")
        generated = os.path.join(out_dir, candidates[0])
    if os.path.abspath(generated) != os.path.abspath(out_path):
        shutil.move(generated, out_path)
    return out_path


def _excel_to_pdf_basic(xlsx_path: str, out_path: str):
    """
    Last-resort: render DataFrame rows to PDF with ReportLab (basic, no formatting).
    """
    import pandas as pd
    from reportlab.pdfgen import canvas
    from reportlab.lib.pagesizes import letter

    df = pd.read_excel(xlsx_path, dtype=str)
    _ensure_dir(Path(out_path).parent.as_posix())
    c = canvas.Canvas(out_path, pagesize=letter)
    y = 750
    for _, row in df.iterrows():
        if y < 50:
            c.showPage()
            y = 750
        row_text = " | ".join([str(v) if v is not None else "" for v in row.values])
        c.drawString(50, y, row_text[:2000])
        y -= 18
    c.save()
    return out_path


def excel_to_pdf(xlsx_path: str, out_path: str):
    """
    Excel -> PDF:
    1) Windows COM (best fidelity)
    2) LibreOffice headless
    3) ReportLab basic fallback
    """
    _ensure_dir(Path(out_path).parent.as_posix())
    if os.name == "nt":
        try:
            return _excel_to_pdf_windows_com(xlsx_path, out_path)
        except Exception:
            try:
                return _office_to_pdf_libreoffice_any(xlsx_path, out_path)
            except Exception:
                return _excel_to_pdf_basic(xlsx_path, out_path)
    else:
        try:
            return _office_to_pdf_libreoffice_any(xlsx_path, out_path)
        except Exception:
            return _excel_to_pdf_basic(xlsx_path, out_path)

def remove_document_metadata(doc_path: str, out_path: str):
    """Remove metadata from document"""
    ext = Path(doc_path).suffix.lower()
    if ext in ['.docx']:
        return _remove_word_metadata(doc_path, out_path)
    elif ext in ['.xlsx']:
        return _remove_excel_metadata(doc_path, out_path)
    else:
        shutil.copy2(doc_path, out_path)
        return out_path

def _remove_word_metadata(doc_path: str, out_path: str):
    from docx import Document
    doc = Document(doc_path)
    props = doc.core_properties
    props.author = ""
    props.last_modified_by = ""
    props.created = None
    props.modified = None
    props.comments = ""
    doc.save(out_path)
    return out_path


# ============== PowerPoint -> PDF ==============

def _ppt_to_pdf_windows_com(ppt_path: str, out_path: str):
    """
    Windows: Use PowerPoint COM automation
    """
    import win32com.client  # pip install pywin32
    ppSaveAsPDF = 32
    app = win32com.client.DispatchEx("PowerPoint.Application")
    app.Visible = True
    pres = app.Presentations.Open(os.path.abspath(ppt_path), WithWindow=False)
    try:
        pres.SaveAs(os.path.abspath(out_path), ppSaveAsPDF)
    finally:
        pres.Close()
        app.Quit()
    return out_path


def _ppt_to_pdf_text_fallback(ppt_path: str, out_path: str):
    """
    Last-resort: extract text via python-pptx and render using ReportLab (no layout).
    """
    from pptx import Presentation
    from reportlab.pdfgen import canvas
    from reportlab.lib.pagesizes import letter

    prs = Presentation(ppt_path)
    _ensure_dir(Path(out_path).parent.as_posix())
    c = canvas.Canvas(out_path, pagesize=letter)
    y = 750
    for slide in prs.slides:
        # New page per slide
        c.showPage()
        y = 750
        for shape in slide.shapes:
            if hasattr(shape, "text") and shape.text:
                if y < 50:
                    c.showPage()
                    y = 750
                c.drawString(50, y, shape.text[:2000])
                y -= 18
    c.save()
    return out_path


def powerpoint_to_pdf(ppt_path: str, out_path: str):
    """
    PowerPoint -> PDF:
    1) Windows COM (PowerPoint)
    2) LibreOffice headless
    3) Text-only fallback
    """
    _ensure_dir(Path(out_path).parent.as_posix())
    if os.name == "nt":
        try:
            return _ppt_to_pdf_windows_com(ppt_path, out_path)
        except Exception:
            try:
                return _office_to_pdf_libreoffice_any(ppt_path, out_path)
            except Exception:
                return _ppt_to_pdf_text_fallback(ppt_path, out_path)
    else:
        try:
            return _office_to_pdf_libreoffice_any(ppt_path, out_path)
        except Exception:
            return _ppt_to_pdf_text_fallback(ppt_path, out_path)


# ============== Class wrapper (backward compatible) ==============

class DocumentConverters:
    """
    Backward-compatible wrapper producing files into a working directory.
    Prefer calling functional API with explicit out_path.
    """
    def __init__(self, output_dir: Optional[str] = None):
        if output_dir:
            self.output_dir = output_dir
            _ensure_dir(self.output_dir)
            self._owns_dir = False
        else:
            self.output_dir = tempfile.mkdtemp(prefix="docs_conv_")
            self._owns_dir = True

    def __del__(self):
        # Keep outputs by default; change if you want auto-clean
        pass

    def pdf_to_word(self, file_path: str) -> str:
        out = os.path.join(self.output_dir, _unique_name("converted", ".docx"))
        return pdf_to_docx(file_path, out)

    def word_to_pdf(self, file_path: str) -> str:
        out = os.path.join(self.output_dir, _unique_name("converted", ".pdf"))
        return docx_to_pdf(file_path, out)

    def pdf_to_excel(self, file_path: str) -> str:
        out = os.path.join(self.output_dir, _unique_name("converted", ".xlsx"))
        return pdf_to_excel(file_path, out)

    def excel_to_pdf(self, file_path: str) -> str:
        out = os.path.join(self.output_dir, _unique_name("converted", ".pdf"))
        return excel_to_pdf(file_path, out)

    def powerpoint_to_pdf(self, file_path: str) -> str:
        out = os.path.join(self.output_dir, _unique_name("converted", ".pdf"))
        return powerpoint_to_pdf(file_path, out)
